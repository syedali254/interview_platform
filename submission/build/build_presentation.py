"""Build the CMP7200 viva voce presentation.

    python build_viva.py

Marking weights for this assessment drive the structure:
    Presentation — visual      15%
    Presentation — verbal      20%
    Critical evaluation        40%   <- the largest single component
    Q&A / discussion           25%

So the deck spends its middle on findings, limitations and reflection rather
than on a tour of the system. Slides carry figures and short claims; the
argument lives in the speaker notes, which are written to be spoken rather
than read aloud verbatim.
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
ROOT = HERE.parent.parent
FIGURES = HERE / "figures_png"
EXP_FIGURES = ROOT / "InterviewAI" / "experiments" / "figures"
STATS = ROOT / "InterviewAI" / "experiments" / "results" / "statistics.json"
OUTPUT = HERE.parent / "CMP7200_Viva_Presentation.pptx"

# ── Palette, matching the dissertation figures ───────────────────────────
INK = RGBColor(0x1B, 0x27, 0x33)
MUTED = RGBColor(0x5A, 0x6B, 0x7B)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
WASH = RGBColor(0xF4, 0xF7, 0xFA)
BLUE = RGBColor(0x1F, 0x4E, 0x79)
GREEN = RGBColor(0x1E, 0x6B, 0x4F)
RUST = RGBColor(0x9A, 0x44, 0x15)
PURPLE = RGBColor(0x4C, 0x3A, 0x8C)
ACCENT = RGBColor(0xB4, 0x53, 0x0F)

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.72)
FONT = "Calibri"


def _stats():
    return json.loads(STATS.read_text(encoding="utf-8")) if STATS.exists() else {}


AUTHOR = "Abdul Wahab"


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = W, H
        self.blank = self.prs.slide_layouts[6]
        self.n = 0

    # ── Slide scaffolding ────────────────────────────────────────────────

    def slide(self, title=None, eyebrow=None, *, notes="", accent=BLUE,
              number=True):
        s = self.prs.slides.add_slide(self.blank)
        self.n += 1
        top = MARGIN

        if eyebrow:
            self.text(s, eyebrow.upper(), MARGIN, top, W - 2 * MARGIN, Inches(0.3),
                      size=12, colour=accent, bold=True, space=1.5)
            top += Inches(0.36)

        if title:
            self.text(s, title, MARGIN, top, W - 2 * MARGIN, Inches(0.72),
                      size=30, colour=INK, bold=True)
            top += Inches(0.86)
            # Accent rule under the title
            self.rule(s, MARGIN, top - Inches(0.10), Inches(1.5), accent)
            top += Inches(0.16)

        if number:
            self.text(s, str(self.n), W - MARGIN - Inches(0.6),
                      H - Inches(0.52), Inches(0.5), Inches(0.3),
                      size=11, colour=MUTED, align=PP_ALIGN.RIGHT)

        if notes:
            s.notes_slide.notes_text_frame.text = notes.strip()

        s._content_top = top
        return s

    def text(self, slide, body, left, top, width, height, *, size=16,
             colour=INK, bold=False, italic=False, align=PP_ALIGN.LEFT,
             space=0, line=1.25, font=FONT):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        for i, line_text in enumerate(str(body).split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.line_spacing = line
            r = p.add_run()
            r.text = line_text
            f = r.font
            f.name, f.size, f.bold, f.italic = font, Pt(size), bold, italic
            f.color.rgb = colour
            if space:
                r.font._rPr.set("spc", str(int(space * 100)))
        return box

    def bullets(self, slide, items, left, top, width, *, size=17, gap=0.52,
                colour=INK, marker_colour=BLUE, lead_colour=None):
        y = top
        for item in items:
            lead, rest = (item if isinstance(item, tuple) else (None, item))
            self.rule(slide, left, y + Inches(0.13), Inches(0.16), marker_colour,
                      thickness=Inches(0.055))
            box = slide.shapes.add_textbox(left + Inches(0.34), y,
                                           width - Inches(0.34), Inches(gap))
            tf = box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.line_spacing = 1.22
            if lead:
                r = p.add_run()
                r.text = lead
                r.font.name, r.font.size, r.font.bold = FONT, Pt(size), True
                r.font.color.rgb = lead_colour or marker_colour
            r = p.add_run()
            r.text = rest
            r.font.name, r.font.size = FONT, Pt(size)
            r.font.color.rgb = colour
            lines = max(1, len(rest) // int(width / Inches(1) * 11) + 1)
            y += Inches(gap * lines * 0.62 + 0.16)
        return y

    def rule(self, slide, left, top, width, colour, thickness=Inches(0.045)):
        from pptx.enum.shapes import MSO_SHAPE
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thickness)
        shp.fill.solid()
        shp.fill.fore_color.rgb = colour
        shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def panel(self, slide, left, top, width, height, colour=WASH, line=None):
        from pptx.enum.shapes import MSO_SHAPE
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shp.fill.solid()
        shp.fill.fore_color.rgb = colour
        if line:
            shp.line.color.rgb = line
            shp.line.width = Pt(1.1)
        else:
            shp.line.fill.background()
        shp.shadow.inherit = False
        shp.adjustments[0] = 0.06
        return shp

    def figure(self, slide, name, *, experiments=False, top=None, height=None,
               left=None, width=None):
        path = (EXP_FIGURES if experiments else FIGURES) / f"{name}.png"
        if not path.exists():
            raise FileNotFoundError(path)
        from PIL import Image
        iw, ih = Image.open(path).size
        ratio = ih / iw
        if height is not None:
            w = int(height / ratio)
            h = height
        else:
            w = width or int(W - 2 * MARGIN)
            h = int(w * ratio)
        l = left if left is not None else int((W - w) / 2)
        t = top if top is not None else int((H - h) / 2)
        slide.shapes.add_picture(str(path), l, t, w, h)

    def stat(self, slide, left, top, width, value, label, colour=BLUE,
             height=Inches(1.5), value_size=40, label_size=13):
        self.panel(slide, left, top, width, height)
        self.text(slide, value, left, top + Inches(0.22), width, Inches(0.7),
                  size=value_size, colour=colour, bold=True, align=PP_ALIGN.CENTER)
        self.text(slide, label, left + Inches(0.16), top + Inches(0.92),
                  width - Inches(0.32), Inches(0.5), size=label_size,
                  colour=MUTED, align=PP_ALIGN.CENTER, line=1.15)

    def save(self):
        # Without this, PowerPoint credits the python-pptx author under
        # File > Properties rather than the person who wrote the deck.
        props = self.prs.core_properties
        props.author = AUTHOR
        props.last_modified_by = AUTHOR
        props.title = "An Explainable Multi-Agent AI Interview Platform"
        self.prs.save(OUTPUT)
        return OUTPUT


# ═════════════════════════════════════════════════════════════════════════

def ensure_figures():
    """Draw the architecture diagrams if they are not on disk yet."""
    expected = [f"fig{n:02d}" for n in range(1, 12)]
    have = {p.stem[:5] for p in FIGURES.glob("*.png")} if FIGURES.exists() else set()
    if not all(e in have for e in expected):
        print("  Rendering architecture figures...")
        import figures
        figures.main()


def module_slide(d, title, eyebrow, rows, footer, notes, accent=BLUE):
    """One slide of modules: what each does, and what it runs on.

    Three columns, because those are the three questions asked of every part:
    which module is it, what does it do, and what was it built with.
    """
    s = d.slide(title, eyebrow=eyebrow, accent=accent, notes=notes)
    y = s._content_top + Inches(0.16)
    d.text(s, "MODULE", MARGIN, y, Inches(1.2), Inches(0.26),
           size=11, colour=MUTED, bold=True, space=1.2)
    d.text(s, "WHAT IT DOES", MARGIN + Inches(1.5), y, Inches(5.0), Inches(0.26),
           size=11, colour=MUTED, bold=True, space=1.2)
    d.text(s, "BUILT WITH", MARGIN + Inches(8.4), y, Inches(3.4), Inches(0.26),
           size=11, colour=MUTED, bold=True, space=1.2)
    y += Inches(0.42)
    d.rule(s, MARGIN, y - Inches(0.08), W - 2 * MARGIN, WASH, thickness=Inches(0.014))

    for tag, name, does, tech, col in rows:
        d.text(s, tag, MARGIN, y, Inches(1.3), Inches(0.34),
               size=19, colour=col, bold=True)
        d.text(s, name, MARGIN + Inches(1.5), y, Inches(6.6), Inches(0.32),
               size=17, colour=INK, bold=True)
        d.text(s, does, MARGIN + Inches(1.5), y + Inches(0.34), Inches(6.6),
               Inches(0.62), size=14.5, colour=MUTED, line=1.26)
        d.text(s, tech, MARGIN + Inches(8.4), y + Inches(0.02), Inches(3.5),
               Inches(0.8), size=14.5, colour=col, bold=True, line=1.26)
        y += Inches(1.06)

    d.panel(s, MARGIN, Inches(6.12), W - 2 * MARGIN, Inches(0.7), WASH)
    d.text(s, footer, MARGIN + Inches(0.42), Inches(6.31),
           W - 2 * MARGIN - Inches(0.84), Inches(0.44),
           size=16, colour=INK, bold=True)
    return s


def build():
    """Ten slides: problem, solution, the modules, results, reflection.

    Four of the ten explain the modules, because that is what the artefact is.
    Each module gets what it does and what it was built with on the same line,
    so the technology question is answered where it is asked rather than on a
    separate slide nobody can map back.
    """
    ensure_figures()
    st = _stats()
    e1 = st.get("e1_discriminant_validity", {})
    e2 = st.get("e2_positional_bias", {})
    e4 = st.get("e4_criterion_independence", {})
    lv = e1.get("by_level", {})
    d = Deck()

    # ── 1. Title ─────────────────────────────────────────────────────────
    s = d.slide(number=False, notes="""
Good morning. My project is an AI system that interviews a job candidate, marks
their answers, and explains how it reached every mark.

The one sentence to hold on to: most automated interview tools give you a score
and no explanation. Mine measures how far its own scores can be trusted, and
says so when they cannot.
""")
    d.rule(s, MARGIN, Inches(1.55), Inches(2.2), ACCENT, thickness=Inches(0.06))
    d.text(s, "An AI Interview Platform\nThat Explains Its Own Marking",
           MARGIN, Inches(1.95), W - 2 * MARGIN, Inches(2.0),
           size=40, colour=INK, bold=True, line=1.14)
    d.text(s, "It interviews a candidate, marks the answers, and shows how it\n"
              "reached every mark \u2014 including how sure it is about each one.",
           MARGIN, Inches(3.9), W - 2 * MARGIN, Inches(1.0),
           size=19, colour=MUTED, line=1.3)
    d.rule(s, MARGIN, Inches(5.3), W - 2 * MARGIN, WASH, thickness=Inches(0.02))
    d.text(s, "Abdul Wahab", MARGIN, Inches(5.6), W - 2 * MARGIN, Inches(0.4),
           size=17, colour=INK, bold=True)
    d.text(s, "Student Number: [STUDENT NUMBER]    \u00b7    CMP7200 Individual "
              "Master's Project    \u00b7    Birmingham City University    \u00b7    2025-26",
           MARGIN, Inches(6.05), W - 2 * MARGIN, Inches(0.4), size=13, colour=MUTED)

    # ── 2. Problem statement ─────────────────────────────────────────────
    s = d.slide("Companies score interviews with AI. Candidates get a number "
                "and no reason.",
                eyebrow="The problem", accent=RUST, notes="""
Companies now use AI to interview and score people at a scale no human team
could manage. It is fast, and it asks everyone the same questions.

The trouble is what comes out: a number, and nothing else. No reason, and no way
to argue with it.

This has caused real trouble. In 2019 a privacy group formally complained about
HireVue scoring candidates from their faces. HireVue dropped that in 2021, but
the deeper problem stayed.

It also costs companies money. Candidates who think a process was unfair are
less likely to accept the job, which cancels out the time saved.

And the law has caught up. The EU AI Act now treats hiring software as high
risk. Notice what it asks for: not a perfect model, which nobody can build, but
that problems be tested for, written down, and checked by a person. That is a
list of requirements, and it is the one I built to.

So my question is not whether AI can mark an answer. It plainly can. It is
whether it can be honest about how far each mark should be trusted.
""")
    y = s._content_top + Inches(0.24)
    for lead, body, col in (
        ("It has caused real trouble.  ",
         "In 2019 a privacy group complained about HireVue scoring candidates "
         "from their faces. The feature went in 2021. The problem stayed.", RUST),
        ("It costs companies money.  ",
         "Candidates who think a process was unfair are less likely to accept "
         "the job \u2014 cancelling out the time saved.", ACCENT),
        ("The law has caught up.  ",
         "The EU AI Act treats hiring software as high risk. It asks that "
         "problems be tested for, written down, and checked by a person.", BLUE),
    ):
        d.rule(s, MARGIN, y + Inches(0.17), Inches(0.2), col, thickness=Inches(0.06))
        d.text(s, lead + body, MARGIN + Inches(0.42), y, W - 2 * MARGIN - Inches(0.5),
               Inches(0.95), size=18, colour=INK, line=1.32)
        y += Inches(1.15)

    d.panel(s, MARGIN, Inches(5.55), W - 2 * MARGIN, Inches(1.2), WASH)
    d.text(s, "The question is not \u201ccan AI mark an interview answer?\u201d "
              "It clearly can.\nIt is \u201ccan it be honest about how much to trust "
              "each mark?\u201d",
           MARGIN + Inches(0.45), Inches(5.8), W - 2 * MARGIN - Inches(0.9),
           Inches(0.85), size=19, colour=INK, bold=True, line=1.35)

    # ── 3. Solution ──────────────────────────────────────────────────────
    s = d.slide("How it works, start to finish", eyebrow="The solution", notes="""
Here is the whole thing, left to right, in four stages.

Prepare. A CV and a job advert go in. The system maps the skills onto ESCO, the
official EU skills list, so when it says a skill is missing, that skill is a real
published thing and not something the AI invented. It writes the interview and
puts the missing skills first, because every interview runs out of time and a
question at the end never gets asked.

Interview. The candidate picks speaking or typing. Both produce exactly the same
transcript, so nothing after this point knows which one ran.

Mark. Greetings and goodbyes are thrown away so they cannot drag the average up.
Every real answer is marked twice. Timing and tab switches get checked.

Report. Everything is combined into one score and a recommendation, with the
working shown.

Thirteen parts do this, and the next four slides go through them.
""")
    d.figure(s, "fig01_architecture", top=s._content_top + Inches(0.1),
             height=Inches(4.0))
    y = Inches(5.98)
    for i, (lead, body, col) in enumerate((
        ("1. Prepare  ", "read both documents, find the missing skills, write "
                         "the questions", BLUE),
        ("2. Interview  ", "by speaking or typing \u2014 both give the same "
                           "transcript", GREEN),
        ("3. Mark  ", "score every answer twice, check the session looked "
                      "normal", RUST),
        ("4. Report  ", "combine into a recommendation, and show the working",
         PURPLE))):
        x = MARGIN + i * ((W - 2 * MARGIN) / 4)
        d.text(s, lead, x, y, Inches(1.7), Inches(0.3), size=14, colour=col, bold=True)
        d.text(s, body, x, y + Inches(0.3), (W - 2 * MARGIN) / 4 - Inches(0.25),
               Inches(0.8), size=12, colour=MUTED, line=1.22)
    d.text(s, "Thirteen parts, four stages \u2014 the next four slides go through them.",
           MARGIN, Inches(6.98), W - 2 * MARGIN, Inches(0.34), size=15,
           colour=INK, bold=True)

    # ── 4. Modules: preparing the interview ──────────────────────────────
    module_slide(d, "Preparing the interview", "Modules 1 to 4", [
        ("M1", "Reads the CV",
         "Pulls out skills, jobs, education and projects from an uploaded PDF "
         "or pasted text. Repairs broken replies rather than failing.",
         "Google Gemini\n3.6 Flash\n+ PyMuPDF", BLUE),
        ("M2", "Reads the job advert",
         "Same treatment for the role: required skills, nice-to-have skills, "
         "seniority and domain.",
         "Google Gemini\n3.6 Flash", BLUE),
        ("M3", "Compares them and finds the gaps",
         "Maps both lists onto ESCO, the EU skills standard, through four "
         "increasingly forgiving matching steps, then reports what is missing.",
         "NetworkX\n+ ESCO v1.1.1\n(1,201 skills)", BLUE),
        ("M4", "Writes the questions",
         "Produces openers, technical questions per skill, behavioural "
         "questions and closers \u2014 then reorders so missing skills come first.",
         "Google Gemini\n3.6 Flash", BLUE),
    ], "Because ESCO is a published standard, \u201cthis skill is missing\u201d means "
       "something a candidate could check.", """
The first four parts turn two documents into a plan.

M1 and M2 read the CV and the job advert. Both use Gemini, asked to return a
fixed structure. Models sometimes return broken data, so the client repairs it
rather than failing the request \u2014 unglamorous, but it is the difference between
a demo that works and one that dies in front of you.

M3 is the core of the project. It maps both skill lists onto ESCO, the EU's
published skills taxonomy, using four matching steps that get gradually more
forgiving. This is deliberately not machine learning: ESCO is already a curated
hierarchy, and learning over it would add opacity while removing the one
property that makes it defensible \u2014 that every skill traces to a published
standard.

M4 writes the interview, then reorders it so the missing skills get asked first.
That reordering is Objective 2, and it matters because every interview has a
time budget and a question at the end never gets asked.
""", accent=BLUE)

    # ── 5. Modules: running the interview ────────────────────────────────
    module_slide(d, "Running the interview", "Modules 5, 5t and 6a", [
        ("M5", "Runs the spoken interview",
         "Shows each question, speaks it aloud, listens to the answer, decides "
         "whether to dig deeper, and saves the transcript.",
         "LiveKit\n+ Deepgram Nova-3\n+ ElevenLabs", GREEN),
        ("M5t", "Runs the typed interview",
         "The same interview answered in a chat box. Produces an identical "
         "transcript, so nothing afterwards knows which mode ran.",
         "FastAPI", GREEN),
        ("M6a", "Tracks how each skill is going",
         "Keeps a live note on every skill \u2014 not asked, asked, answered well, "
         "answered badly \u2014 which decides whether a follow-up is needed.",
         "Plain Python\n(state machine)", GREEN),
    ], "M5 and M5t are the same interview with a different way of answering \u2014 "
       "which is why it is thirteen parts, not fourteen.", """
The interview itself.

M5 is the talking interviewer. LiveKit carries the audio call, Deepgram turns
speech into text, and ElevenLabs is the voice \u2014 with Deepgram's own voice as a
backup if that runs out of credit. It shows each question on screen before
speaking it, so a candidate who mishears can read it.

Two things I got wrong here and had to fix. The system checks the voice provider
actually returns audio before starting, because an exhausted account accepts the
connection and returns silence, which looks exactly like success. And the
turn-taking timings had to be slowed right down: the defaults are tuned for
chat, and an interview answer contains long pauses for thought.

M5t is the same interview, typed. It matters more than it sounds: it produces an
identical transcript, so every part after this point is shared. It also lets you
demonstrate the whole system without any audio hardware at all.

M6a keeps the running verdict on each skill that decides whether to follow up.
""", accent=GREEN)

    # ── 6. Modules: watching the interview ───────────────────────────────
    module_slide(d, "Watching the interview", "Modules 7, 8 and 10", [
        ("M7", "Watches attention",
         "Finds 478 points on the face to work out where the head is pointing "
         "and whether the eyes are on the screen.",
         "MediaPipe\nFace Landmarker", GREEN),
        ("M8", "Watches posture",
         "Finds 33 points on the upper body to check whether the shoulders are "
         "level and the person is upright rather than slouched.",
         "MediaPipe\nPose Landmarker", GREEN),
        ("M10", "Listens to tone of voice",
         "Measures the sound, not the words: how steady the pitch is, how loud, "
         "and how much is speech rather than pauses.",
         "Web Audio API", GREEN),
    ], "All three run inside the candidate's own browser. Video and audio never "
       "leave their computer \u2014 only the resulting numbers are sent.", """
Three parts watch how the interview went rather than what was said.

M7 finds 478 points on the face and works out where the head is pointing and
whether the eyes are on the screen. M8 finds 33 points on the upper body and
checks posture. M10 listens to the sound of the voice \u2014 pitch steadiness,
loudness, how much of the time is actual speech.

The line at the bottom is the one I would emphasise. All three run inside the
candidate's own browser using models downloaded once. The video and the audio
never leave their machine; only the resulting numbers are sent. That is what
makes the privacy claim real rather than a promise.

Two limits I would state before you ask. None of this reads emotion \u2014 I removed
a facial emotion classifier precisely because it could not be explained, and my
Chapter 2 criticises HireVue for exactly that. And all three together are worth
fifteen points out of a hundred in the final score, against fifty for what the
candidate actually said. Nobody is marked on their face.
""", accent=GREEN)

    # ── 7. Modules: marking and reporting ────────────────────────────────
    module_slide(d, "Marking and reporting", "Modules 6, 9, 11 and 12", [
        ("M6", "Marks every answer \u2014 twice",
         "Scores each answer on accuracy, completeness, clarity and relevance, "
         "with the criteria in a different order each time. The gap between "
         "the two marks says how far to trust it.",
         "Gemini as judge\n(permuted order)", RUST),
        ("M9", "Checks for cheating",
         "Looks at tab switches, answers arriving too fast or too slow, and "
         "answers that are suspiciously uniform in length.",
         "scikit-learn\nIsolation Forest\n(trains itself)", RUST),
        ("M11", "Adds everything up",
         "Answers 50%, skills covered 20%, session integrity 15%, delivery "
         "15% \u2014 with every contribution shown, not hidden.",
         "Plain Python\n(fixed weights)", PURPLE),
        ("M12", "Writes the report",
         "Each answer and mark, how consistent the marking was, the integrity "
         "verdict and the recommendation.",
         "Plain Python", PURPLE),
    ], "M6 is the research contribution: mark twice, and keep the disagreement "
       "as a measure of how far to trust the mark.", """
The marking, and this is where the research contribution sits.

M6 is the important one. Research shows that if you change the order you present
the marking criteria in, an AI judge gives a different score. Most people cite
that as a reason not to trust AI marking. I read it as an instruction: mark every
answer twice, with the criteria in a different order each time.

Averaging cancels the effect of the order. But the average is not the interesting
part \u2014 the gap between the two marks is. Eighty-two then eighty-one means it is
sure. Seventy-one then forty-five means it is not, and the average of fifty-eight
hides that completely. So the gap is kept, and a big gap sends that answer to a
human instead of reporting a confident score.

M9 checks the session behaved normally. It trains its own model of normal on
first run, because labelled examples of cheating are not something anyone can
ethically collect. It always names the behaviour that caused a flag.

M11 combines everything on fixed published weights \u2014 half the score is what the
candidate actually said. M12 writes it up with the working shown.
""", accent=RUST)

    # ── 8. Results ───────────────────────────────────────────────────────
    s = d.slide("What testing found \u2014 two faults in my own system",
                eyebrow="Results and impact", accent=RUST, notes="""
I wrote 18 answers to a known quality \u2014 deliberately weak, average and excellent
\u2014 and ran five experiments.

The good news. It almost never puts a worse answer above a better one: 0.92 out
of 1. As a way of ranking candidates against each other, it works.

The first fault. Weak answers average 53. Deliberately average answers average
92.8. Excellent ones 98.2. And the system calls anything 70 or above strong. So
average and excellent get the same verdict. These marks are usable for comparing
two candidates; they are not usable for deciding whether one candidate is good
enough on their own. My report says exactly that.

The second fault. I ask for four separate criteria. They move together at 0.85.
It forms one overall impression and spreads it across all four. Telling the model
not to do that did not stop it \u2014 and it weakens a claim I made in my own design
chapter, that the breakdown explains which part of an answer fell short.

And an honest non-result: the safety net never fired, because the two marks were
only 2.2 points apart on average. I report that as a null result rather than
dressing it up.

The impact line is the one to land. Both faults are in my own system and my own
testing found them. That is the entire argument for building it this way.
""")
    y = s._content_top + Inches(0.1)
    d.stat(s, MARGIN, y, Inches(2.75), "0.92",
           "how well it ranks answers\n(1.0 would be perfect)", GREEN,
           height=Inches(1.4), value_size=36)
    d.stat(s, MARGIN + Inches(3.0), y, Inches(2.75),
           f"{lv.get('medium', {}).get('mean', 0):.1f}",
           "what deliberately average\nanswers scored", RUST,
           height=Inches(1.4), value_size=36)
    d.stat(s, MARGIN + Inches(6.0), y, Inches(2.75),
           f"{e4.get('mean_inter_criterion_r', 0):.2f}",
           "how closely the four\ncriteria move together", RUST,
           height=Inches(1.4), value_size=36)
    d.stat(s, MARGIN + Inches(9.0), y, Inches(2.85),
           f"{e2.get('mean_absolute_spread', 0):.1f}",
           "average gap between\nthe two markings", BLUE,
           height=Inches(1.4), value_size=36)
    y += Inches(1.7)
    for lead, body, col in (
        ("It ranks well.  ",
         "It almost never puts a worse answer above a better one.", GREEN),
        ("Fault 1 \u2014 too generous.  ",
         "The pass mark for \u201cstrong\u201d is 70. Deliberately average answers scored "
         "92.8, so average and excellent get the same verdict.", RUST),
        ("Fault 2 \u2014 the four criteria are not separate.  ",
         "It forms one impression and spreads it across all four, so the "
         "breakdown explains less than I had claimed.", RUST),
    ):
        d.rule(s, MARGIN, y + Inches(0.15), Inches(0.18), col, thickness=Inches(0.06))
        d.text(s, lead + body, MARGIN + Inches(0.4), y, W - 2 * MARGIN - Inches(0.5),
               Inches(0.8), size=17, colour=INK, line=1.28)
        y += Inches(0.86)
    d.panel(s, MARGIN, Inches(6.15), W - 2 * MARGIN, Inches(0.72), WASH)
    d.text(s, "Both faults are in my own system, and my own testing found them. "
              "That is the argument for building it this way.",
           MARGIN + Inches(0.42), Inches(6.35), W - 2 * MARGIN - Inches(0.84),
           Inches(0.45), size=17, colour=INK, bold=True)

    # ── 9. Reflection ────────────────────────────────────────────────────
    s = d.slide("What I changed, and what this project cannot claim",
                eyebrow="Looking back honestly", accent=PURPLE, notes="""
Where I departed from my proposal.

My proposal's headline idea was a second marker \u2014 a trained model to check the AI
judge. I built it, measured it, and removed it. Its training answers came from
the same AI it was meant to check, so the two agreeing tells you nothing. A test
showed it was really keying on answer length. And a second unchecked marker gives
the appearance of safety, not safety.

Removing it late felt like going backwards \u2014 it was the part that most looked
like machine learning research. What changed my mind was realising I could not
answer the obvious question: why does that agreement mean anything? Writing it up
as a finding turned the weakest part of the project into one I can defend.

I also replaced a planned emotion detector with a measure of how someone speaks.
Every part of that can be inspected; an emotion label from a black box cannot.

What I cannot claim. Eighteen answers is a small test set, and they were written
rather than spoken. No human markers to compare against. No bias testing across
groups, because that needs exactly the personal data I chose not to collect. And
it runs one interview at a time.

Against the EU AI Act: explaining itself and keeping a human in charge, it does
well. Record keeping and bias testing, it does not.
""")
    y = s._content_top + Inches(0.08)
    d.text(s, "What I changed, and why", MARGIN, y, Inches(5.8), Inches(0.34),
           size=18, colour=PURPLE, bold=True)
    d.text(s, "What it cannot claim", MARGIN + Inches(6.6), y, Inches(5.5),
           Inches(0.34), size=18, colour=RUST, bold=True)
    y += Inches(0.52)
    left = [
        ("Removed the second marker. ",
         "Its training answers came from the same AI it was meant to check, so "
         "agreement proved nothing. A test showed it measuring length."),
        ("Replaced the emotion detector. ",
         "Now measures how someone speaks. Every part can be inspected; an "
         "emotion label cannot."),
        ("Added the double marking. ",
         "Not in the proposal at all. It became the main contribution."),
    ]
    right = [
        ("18 answers, written not spoken. ", "A small test set."),
        ("No human markers to compare against. ", "So no outside check."),
        ("No bias testing across groups. ",
         "It needs the personal data I chose not to collect."),
        ("One interview at a time. ", "A prototype, not a product."),
    ]
    yl = y
    for lead, rest in left:
        d.rule(s, MARGIN, yl + Inches(0.14), Inches(0.16), PURPLE, thickness=Inches(0.055))
        d.text(s, lead + rest, MARGIN + Inches(0.34), yl, Inches(5.7), Inches(1.0),
               size=15, colour=INK, line=1.28)
        yl += Inches(1.02)
    yr = y
    for lead, rest in right:
        d.rule(s, MARGIN + Inches(6.6), yr + Inches(0.14), Inches(0.16), RUST,
               thickness=Inches(0.055))
        d.text(s, lead + rest, MARGIN + Inches(6.94), yr, Inches(5.1), Inches(0.8),
               size=15, colour=INK, line=1.28)
        yr += Inches(0.76)
    d.panel(s, MARGIN, Inches(6.2), W - 2 * MARGIN, Inches(0.68), WASH)
    d.text(s, "Every one of these changes came from a measurement that proved me "
              "wrong about something.",
           MARGIN + Inches(0.42), Inches(6.38), W - 2 * MARGIN - Inches(0.84),
           Inches(0.42), size=17, colour=INK, bold=True)

    # ── 10. Conclusion ───────────────────────────────────────────────────
    s = d.slide("What this project contributes", eyebrow="In closing",
                accent=GREEN, notes="""
Four things.

A working platform that aims its questions using a published skills standard,
interviews by speaking or typing, and produces a report that shows its working.

A clear measurement of how an AI marker behaves in practice \u2014 that it can rank
answers well and still be far too generous. That distinction is where both faults
came from.

A documented negative result: comparing a trained model against the AI that
supplied its training answers proves nothing, and I have the evidence.

And the whole test setup is in the submission, so every number can be produced
again.

The wider point is deliberately modest. Being open about one marker \u2014 what it
did, how sure it was, when a human should step in \u2014 is worth more than adding a
second marker nobody has checked.

And the best evidence for that is what I said earlier: my own testing found two
real faults in my own system. Something built to be examined can be shown to be
wrong, and then fixed. That is worth more than a better-looking score from
something nobody can see inside.

Thank you. I am happy to take questions.
""")
    y = s._content_top + Inches(0.15)
    for lead, body in (
        ("A working platform  ",
         "that aims its questions using a published skills standard, interviews "
         "by speaking or typing, and shows its working."),
        ("A clear measurement  ",
         "of how an AI marker behaves in practice \u2014 it can rank answers well and "
         "still mark them far too generously."),
        ("A documented negative result  ",
         "checking an AI against a model trained on that same AI's answers "
         "proves nothing, and here is the evidence."),
        ("A test setup anyone can rerun  ",
         "included in the submission, so every number here can be produced again."),
    ):
        d.rule(s, MARGIN, y + Inches(0.15), Inches(0.18), GREEN, thickness=Inches(0.06))
        d.text(s, lead + body, MARGIN + Inches(0.4), y, W - 2 * MARGIN - Inches(0.5),
               Inches(0.8), size=17, colour=INK, line=1.3)
        y += Inches(0.88)

    d.panel(s, MARGIN, Inches(5.35), W - 2 * MARGIN, Inches(1.45), WASH)
    d.text(s, "My own testing found two real faults in my own system.",
           MARGIN + Inches(0.45), Inches(5.58), W - 2 * MARGIN - Inches(0.9),
           Inches(0.45), size=19, colour=INK, bold=True)
    d.text(s, "Something built to be examined can be shown to be wrong, and then "
              "fixed. That is worth more than a\nbetter-looking score from "
              "something nobody can see inside.",
           MARGIN + Inches(0.45), Inches(6.05), W - 2 * MARGIN - Inches(0.9),
           Inches(0.7), size=16, colour=MUTED, line=1.3)

    return d.save()


if __name__ == "__main__":
    out = build()
    print(f"Written: {out}")
    from pptx import Presentation as _P
    p = _P(out)
    slides = list(p.slides)
    print(f"  slides        : {len(slides)}")
    print(f"  speaker notes : {sum(1 for s in slides if s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip())}")
    bad = 0
    for i, s in enumerate(slides, 1):
        for sh in s.shapes:
            if sh.left is None:
                continue
            if (sh.left < -10000 or sh.top < -10000
                    or sh.left + (sh.width or 0) > p.slide_width + 10000
                    or sh.top + (sh.height or 0) > p.slide_height + 10000):
                bad += 1
                print(f"  ! slide {i}: bottom={(sh.top + (sh.height or 0)) / 914400:.2f} "
                      f"right={(sh.left + (sh.width or 0)) / 914400:.2f}")
    print(f"bounds problems: {bad or 'none'}")
